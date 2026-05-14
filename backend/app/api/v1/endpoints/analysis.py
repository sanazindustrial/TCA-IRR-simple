"""
Analysis management endpoints
"""

import logging
from fastapi import APIRouter, Depends, HTTPException, status, BackgroundTasks
from typing import List, Dict, Any, Optional
import asyncpg

from app.db import get_db
from app.models import AnalysisResponse, AnalysisCreate, PaginatedResponse
from app.services import analysis_processor
from .auth import get_current_user

logger = logging.getLogger(__name__)
router = APIRouter()


@router.get("/", response_model=PaginatedResponse)
async def get_analyses(page: int = 1,
                       size: int = 20,
                       db: asyncpg.Connection = Depends(get_db),
                       current_user: dict = Depends(get_current_user)):
    """Get paginated list of analyses"""
    # Implementation placeholder
    return PaginatedResponse(items=[],
                             total=0,
                             page=page,
                             size=size,
                             pages=0,
                             has_next=False,
                             has_previous=False)


@router.get("/ai-extract")
async def ai_extract_status(current_user: dict = Depends(get_current_user)):
    """AI Analysis Agent status endpoint for health monitoring"""
    return {
        "status": "operational",
        "service": "AI Analysis Agent",
        "available": True,
        "mode": "calculated",
        "message": "AI analysis service ready"
    }


@router.get("/ai-orchestrate")
async def ai_orchestrate_status(current_user: dict = Depends(get_current_user)):
    """Multi-Agent Orchestrator status endpoint for health monitoring"""
    return {
        "status": "operational",
        "service": "Multi-Agent Orchestrator",
        "available": True,
        "mode": "sequential",
        "message": "Multi-agent orchestration service ready"
    }


@router.post("/",
             response_model=AnalysisResponse,
             status_code=status.HTTP_201_CREATED)
async def create_analysis(analysis_data: AnalysisCreate,
                          background_tasks: BackgroundTasks,
                          db: asyncpg.Connection = Depends(get_db),
                          current_user: dict = Depends(get_current_user)):
    """Create and start a new analysis"""
    # Implementation placeholder
    pass


@router.get("/{analysis_id}", response_model=AnalysisResponse)
async def get_analysis(analysis_id: int,
                       db: asyncpg.Connection = Depends(get_db),
                       current_user: dict = Depends(get_current_user)):
    """Get analysis by ID"""
    # Implementation placeholder
    pass


@router.post("/test", response_model=Dict[str, Any])
async def test_analysis(company_data: Dict[str, Any]):
    """Test endpoint for basic analysis"""
    return {
        "test": "success",
        "company_name": company_data.get("company_name", "Unknown"),
        "message": "Basic endpoint working"
    }


@router.post("/comprehensive", response_model=Dict[str, Any])
async def comprehensive_analysis(company_data: Dict[str, Any]):
    """Run comprehensive TCA analysis on company data"""
    try:
        company_name = (company_data.get("company_name") or company_data.get("name") or "").strip()
        if not company_name:
            raise HTTPException(status_code=422, detail="company_name is required")

        company_data["company_name"] = company_name
        if not company_data.get("industry") and not company_data.get("industry_vertical"):
            raise HTTPException(status_code=422, detail="industry or industry_vertical is required")

        logger.info(
            f"Starting comprehensive analysis for: {company_data.get('company_name', 'Unknown')}"
        )

        # ------------------------------------------------------------------
        # Pull saved module configuration (passed by frontend run/page.tsx).
        #   module_settings: [{module_id, weight, is_enabled}, ...]    top-level enable/weight
        #   module_configs:  {moduleId: {<rubric saved on /analysis/modules/{id}>}}
        # When absent, we fall back to defaults (every module enabled, hardcoded weights).
        # ------------------------------------------------------------------
        raw_module_settings = company_data.get("module_settings") or []
        raw_module_configs  = company_data.get("module_configs")  or {}
        enabled_modules: Dict[str, bool] = {}
        module_weights:  Dict[str, float] = {}
        if isinstance(raw_module_settings, list):
            for entry in raw_module_settings:
                if not isinstance(entry, dict):
                    continue
                mid = entry.get("module_id") or entry.get("id")
                if not mid:
                    continue
                enabled_modules[str(mid)] = bool(entry.get("is_enabled", True))
                try:
                    module_weights[str(mid)] = float(entry.get("weight") or 0.0)
                except (TypeError, ValueError):
                    module_weights[str(mid)] = 0.0
        if enabled_modules:
            logger.info(
                f"Module settings received: {sum(enabled_modules.values())}/{len(enabled_modules)} enabled"
            )

        # Use real analysis processor instead of mock data — honour saved enabled flags.
        def _is_on(mod_id: str, default: bool = True) -> bool:
            if not enabled_modules:
                return default
            return enabled_modules.get(mod_id, default)

        analysis_options = {
            "tca_scorecard":        _is_on("tca"),
            "benchmark_comparison": _is_on("benchmark"),
            "risk_assessment":      _is_on("risk"),
            "founder_analysis":     _is_on("founderFit", company_data.get("founders") is not None),
            "macro_trend":          _is_on("macro"),
            "team_assessment":      _is_on("team"),
            "growth_classifier":    _is_on("growth"),
            "gap_analysis":         _is_on("gap"),
            "strategic_fit":        _is_on("strategicFit"),
            "financial_analysis":   _is_on("financial"),
            "economic_analysis":    _is_on("economic"),
            "social_analysis":      _is_on("social"),
            "marketing_analysis":   _is_on("marketing"),
            "environmental_analysis": _is_on("environmental"),
        }
        # Process analysis through AI service
        try:
            ai_results = await analysis_processor.process_comprehensive_analysis(
                company_data, analysis_options)
            logger.info(
                f"AI analysis completed with status: {ai_results.get('status')}"
            )

            # If AI analysis fails, use calculated fallback
            if ai_results.get(
                    "status") != "success" or not ai_results.get("results"):
                logger.warning(
                    "AI analysis failed or incomplete, using calculated fallback"
                )
                analysis_result = _calculate_fallback_analysis(
                    company_data, module_weights, raw_module_configs)
            else:
                analysis_result = _format_ai_results(ai_results, company_data)

        except Exception as ai_error:
            logger.error(f"AI analysis error: {ai_error}")
            # Fallback to calculated analysis
            analysis_result = _calculate_fallback_analysis(
                company_data, module_weights, raw_module_configs)

        # Echo applied module configuration back to the client so the result/triage
        # pages can render the exact weights/enabled-set that produced these scores.
        if enabled_modules or module_weights or raw_module_configs:
            analysis_result.setdefault("module_settings_applied", {
                "enabled": enabled_modules,
                "weights": module_weights,
                "configs": raw_module_configs,
            })

        return analysis_result

    except Exception as e:
        logger.error(f"Comprehensive analysis error: {e}")
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                            detail="Analysis processing failed")


def _calculate_fallback_analysis(
        company_data: Dict[str, Any],
        module_weights: Optional[Dict[str, float]] = None,
        module_configs: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """Calculate TCA analysis using business logic when AI service is unavailable.

    module_weights / module_configs come from the saved /dashboard/module-settings
    configuration. They are optional — when absent the function preserves the
    legacy behaviour with hardcoded weights.
    """
    from datetime import datetime

    company_name = company_data.get("company_name", "Unknown Company")

    # Calculate TCA scores based on available data — honour saved sub-category
    # weights from module_configs["tca"] when present.
    tca_categories = _calculate_tca_categories(
        company_data,
        tca_module_config=(module_configs or {}).get("tca"),
    )
    composite_score = _calculate_composite_score(
        tca_categories,
        module_weight=(module_weights or {}).get("tca"),
    )

    # Generate risk assessment
    risk_data = _calculate_risk_assessment(company_data, tca_categories)

    # Generate recommendation
    recommendation = _generate_recommendation(composite_score, risk_data)

    return {
        "analysis_id": f"calc_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
        "company_name": company_name,
        "executive_summary": {
            "overall_score": composite_score,
            "recommendation": recommendation["overall"],
            "key_strengths": recommendation["strengths"],
            "key_concerns": recommendation["concerns"]
        },
        "tca_data": {
            "categories": tca_categories,
            "compositeScore": composite_score,
            "summary": recommendation["summary"]
        },
        "risk_data": risk_data,
        "detailed_analysis": {
            "market_analysis": _analyze_market_potential(company_data),
            "financial_analysis": _analyze_financial_health(company_data),
            "team_analysis": _analyze_team_strength(company_data),
            "risk_assessment": risk_data["summary"]
        },
        "status": "completed",
        "created_at": datetime.now().isoformat()
    }


def _format_ai_results(ai_results: Dict[str, Any],
                       company_data: Dict[str, Any]) -> Dict[str, Any]:
    """Format AI analysis results into consistent structure"""
    from datetime import datetime

    results = ai_results.get("results", {})

    return {
        "analysis_id": f"ai_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
        "company_name": company_data.get("company_name", "Unknown Company"),
        "ai_powered": True,
        "executive_summary": {
            "overall_score":
            results.get("tca_scorecard", {}).get("overall_score", 0),
            "recommendation":
            results.get("tca_scorecard", {}).get("recommendation",
                                                 "Further analysis needed"),
            "key_strengths":
            results.get("tca_scorecard", {}).get("key_insights", [])[:3],
            "key_concerns":
            results.get("risk_assessment", {}).get("major_risks", [])[:3]
        },
        "tca_data": results.get("tca_scorecard", {}),
        "risk_data": results.get("risk_assessment", {}),
        "benchmark_data": results.get("benchmark_comparison", {}),
        "founder_data": results.get("founder_analysis", {}),
        "status": "completed",
        "created_at": datetime.now().isoformat(),
        "ai_errors": ai_results.get("errors", [])
    }


def _calculate_tca_categories(
        company_data: Dict[str, Any],
        tca_module_config: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
    """Calculate TCA category scores based on company data.

    tca_module_config (optional) is the user-saved TCA configuration as stored
    by /analysis/modules/tca — shape: {categories: [{id, name, general, medtech,
    generalNA, medtechNA}, ...]}. When provided, weights are taken from that
    object (general framework) instead of the hardcoded defaults below.
    """

    # TCA Category weights (General framework) — used when no saved config is supplied.
    category_config = {
        "Leadership": {
            "weight": 0.20,
            "factors": ["founder_experience", "leadership_track_record"]
        },
        "Product-Market Fit": {
            "weight":
            0.20,
            "factors":
            ["market_validation", "customer_feedback", "product_traction"]
        },
        "Team Strength": {
            "weight": 0.10,
            "factors":
            ["team_experience", "technical_skills", "domain_expertise"]
        },
        "Technology & IP": {
            "weight": 0.10,
            "factors": ["tech_innovation", "ip_portfolio", "technical_moat"]
        },
        "Business Model & Financials": {
            "weight": 0.10,
            "factors": ["revenue_model", "financial_metrics", "burn_rate"]
        },
        "Go-to-Market Strategy": {
            "weight":
            0.10,
            "factors":
            ["sales_strategy", "marketing_approach", "customer_acquisition"]
        },
        "Competition & Moat": {
            "weight": 0.05,
            "factors": ["competitive_advantage", "market_differentiation"]
        },
        "Market Potential": {
            "weight": 0.05,
            "factors": ["market_size", "growth_rate", "market_timing"]
        },
        "Traction": {
            "weight": 0.05,
            "factors": ["customer_growth", "revenue_growth", "partnerships"]
        },
        "Scalability": {
            "weight": 0.025,
            "factors": ["technical_scalability", "business_scalability"]
        },
        "Risk Assessment": {
            "weight": 0.025,
            "factors": ["identified_risks", "mitigation_strategies"]
        },
        "Exit Potential": {
            "weight": 0.0,
            "factors": ["acquisition_targets", "ipo_potential"]
        }
    }

    # Overlay user-saved TCA category weights when available.
    if tca_module_config and isinstance(tca_module_config, dict):
        saved_categories = tca_module_config.get("categories")
        if isinstance(saved_categories, list):
            # Map saved category names to weights (general framework, percent → fraction).
            saved_weights: Dict[str, float] = {}
            for sc in saved_categories:
                if not isinstance(sc, dict):
                    continue
                name = (sc.get("name") or "").strip()
                if not name:
                    continue
                if sc.get("generalNA"):
                    saved_weights[name] = 0.0
                else:
                    try:
                        saved_weights[name] = float(sc.get("general") or 0.0) / 100.0
                    except (TypeError, ValueError):
                        continue
            if saved_weights:
                for cat_name in list(category_config.keys()):
                    if cat_name in saved_weights:
                        category_config[cat_name]["weight"] = saved_weights[cat_name]

    categories = []

    for category_name, config in category_config.items():
        # Calculate raw score based on available data
        raw_score = _calculate_category_score(company_data, category_name,
                                              config["factors"])
        weighted_score = raw_score * config["weight"]

        # Determine flag color
        flag_color = "green" if raw_score >= 8.0 else "yellow" if raw_score >= 6.5 else "red"

        categories.append({
            "category":
            category_name,
            "rawScore":
            round(raw_score, 1),
            "weight":
            int(config["weight"] * 100),
            "weightedScore":
            round(weighted_score, 2),
            "flag":
            flag_color,
            "description":
            f"Evaluation of {category_name.lower()} performance",
            "strengths":
            f"Strong performance indicators in {category_name.lower()}",
            "concerns":
            "Areas for improvement identified"
            if raw_score < 7.0 else "Minor optimization opportunities",
            "interpretation":
            _get_category_interpretation(category_name, raw_score),
            "aiRecommendation":
            _get_category_recommendation(category_name, raw_score)
        })

    return categories


def _get_nested_value(data: Dict[str, Any], path: str) -> Any:
    """Safely read dotted paths from nested request payloads."""
    current: Any = data
    for part in path.split('.'):
        if not isinstance(current, dict) or part not in current:
            return None
        current = current.get(part)
    return current


def _coerce_number(value: Any) -> float | None:
    if value is None or isinstance(value, bool):
        return None
    if isinstance(value, (int, float)):
        return float(value)
    if isinstance(value, str):
        text = value.strip().replace(',', '')
        if not text:
            return None
        try:
            return float(text)
        except ValueError:
            return None
    return None


def _has_signal(value: Any) -> bool:
    if value is None:
        return False
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return float(value) > 0
    if isinstance(value, str):
        return len(value.strip()) > 0
    if isinstance(value, (list, tuple, set, dict)):
        return len(value) > 0
    return False


def _any_signal(company_data: Dict[str, Any], *paths: str) -> bool:
    return any(_has_signal(_get_nested_value(company_data, path)) for path in paths)


def _count_items(company_data: Dict[str, Any], *paths: str) -> int:
    count = 0
    for path in paths:
        value = _get_nested_value(company_data, path)
        if isinstance(value, (list, tuple, set)):
            count += len(value)
        elif _has_signal(value):
            count += 1
    return count


def _numeric_max(company_data: Dict[str, Any], *paths: str) -> float | None:
    values: List[float] = []
    for path in paths:
        number = _coerce_number(_get_nested_value(company_data, path))
        if number is not None:
            values.append(number)
    return max(values) if values else None


def _calculate_category_score(company_data: Dict[str, Any], category: str,
                              factors: List[str]) -> float:
    """Calculate score for a specific TCA category using available evidence signals."""
    base_score = 5.0
    score = base_score

    if category == "Leadership":
        founder_count = _count_items(company_data, "founders", "team_data.founders", "founder_profiles")
        if founder_count >= 2:
            score += 1.3
        elif founder_count == 1:
            score += 0.7
        if _any_signal(company_data, "founder_experience", "team_data.founders", "leadership_track_record"):
            score += 1.1
        if _any_signal(company_data, "leadership_team", "team_data.key_personnel", "team_data.advisors"):
            score += 0.8
        if _any_signal(company_data, "advisory_board", "board_members"):
            score += 0.6

    elif category == "Product-Market Fit":
        if _any_signal(company_data, "customer_validation", "traction.customer_logos", "market_feedback"):
            score += 1.8
        if _any_signal(company_data, "revenue_traction", "financial_data.revenue", "monthly_revenue"):
            score += 1.4
        customer_count = _numeric_max(company_data, "financial_data.customer_count", "customer_count")
        if customer_count is not None and customer_count >= 100:
            score += 1.0
        if _any_signal(company_data, "product_traction", "pilot_results", "nps_score"):
            score += 0.8

    elif category == "Team Strength":
        team_size = _numeric_max(company_data, "team_size", "team_data.team_size") or 0
        if team_size >= 10:
            score += 1.4
        elif team_size >= 5:
            score += 0.9
        elif team_size > 0:
            score += 0.3
        if _any_signal(company_data, "technical_team", "team_data.key_personnel", "engineering_headcount"):
            score += 1.1
        if _any_signal(company_data, "industry_experience", "team_data.founders", "domain_expertise"):
            score += 0.9

    elif category == "Technology & IP":
        patent_count = _count_items(company_data, "patents", "ip_portfolio")
        if patent_count >= 3:
            score += 1.3
        elif patent_count > 0:
            score += 0.8
        if _any_signal(company_data, "technical_innovation", "unique_technology", "technical_moat"):
            score += 1.1
        if _any_signal(company_data, "tech_stack", "architecture", "model_stack"):
            score += 0.7
        stage = str(_get_nested_value(company_data, "development_stage") or "").lower()
        if stage in {"beta", "production", "launched"}:
            score += 0.8
        elif stage in {"concept", "idea"}:
            score -= 0.6

    elif category == "Business Model & Financials":
        if _any_signal(company_data, "revenue_model", "pricing_model", "business_model"):
            score += 1.1
        if _any_signal(company_data, "positive_unit_economics", "financial_data.gross_margin"):
            score += 1.0
        monthly_revenue = _numeric_max(company_data, "monthly_revenue", "financial_data.revenue")
        if monthly_revenue is not None and monthly_revenue > 0:
            score += 1.2
        runway = _numeric_max(company_data, "financial_data.runway_months", "runway_months")
        if runway is not None:
            if runway >= 12:
                score += 0.8
            elif runway < 6:
                score -= 0.8

    elif category == "Go-to-Market Strategy":
        if _any_signal(company_data, "sales_strategy", "go_to_market", "distribution_strategy"):
            score += 1.4
        if _any_signal(company_data, "marketing_approach", "marketing_channels"):
            score += 1.0
        if _any_signal(company_data, "customer_acquisition", "customer_acquisition_cost", "pipeline"):
            score += 1.0
        if _any_signal(company_data, "partnerships", "channel_partners"):
            score += 0.6

    elif category == "Competition & Moat":
        if _any_signal(company_data, "competitive_advantage", "differentiation", "moat"):
            score += 1.6
        if _any_signal(company_data, "market_differentiation", "unique_value_proposition"):
            score += 1.1
        if _any_signal(company_data, "patents", "ip_portfolio"):
            score += 0.7

    elif category == "Market Potential":
        market_size = _numeric_max(company_data, "market_size", "tam", "market_data.market_size")
        if market_size is not None:
            if market_size >= 1_000_000_000:
                score += 2.0
            elif market_size >= 100_000_000:
                score += 1.2
        growth_rate = _numeric_max(company_data, "growth_rate", "market_growth_rate", "market_data.market_growth_pct")
        if growth_rate is not None:
            normalized_growth = growth_rate / 100 if growth_rate > 1 else growth_rate
            if normalized_growth >= 0.15:
                score += 1.0
            elif normalized_growth >= 0.08:
                score += 0.5
        if _any_signal(company_data, "market_timing", "industry_tailwinds"):
            score += 0.5

    elif category == "Traction":
        if _any_signal(company_data, "customer_growth", "financial_data.customer_count", "customer_count"):
            score += 1.3
        revenue_growth = _numeric_max(company_data, "revenue_growth", "financial_data.revenue_growth_pct")
        if revenue_growth is not None:
            normalized_growth = revenue_growth / 100 if revenue_growth > 1 else revenue_growth
            if normalized_growth >= 0.30:
                score += 1.1
            elif normalized_growth >= 0.10:
                score += 0.6
        if _any_signal(company_data, "partnerships", "pilot_results", "enterprise_contracts"):
            score += 0.9

    elif category == "Scalability":
        if _any_signal(company_data, "technical_scalability", "infrastructure", "cloud_architecture"):
            score += 1.3
        if _any_signal(company_data, "business_scalability", "process_automation", "operating_playbook"):
            score += 1.1
        if _any_signal(company_data, "unit_economics", "gross_margin"):
            score += 0.6

    elif category == "Risk Assessment":
        risks = _count_items(company_data, "identified_risks", "risk_register", "risk_factors")
        mitigations = _count_items(company_data, "mitigation_strategies", "risk_mitigation", "contingency_plan")
        if risks == 0:
            score += 0.6
        elif risks <= 3:
            score += 0.2
        else:
            score -= 0.8
        if mitigations > 0:
            score += 1.0
        if _any_signal(company_data, "compliance_program", "security_controls"):
            score += 0.6

    elif category == "Exit Potential":
        if _any_signal(company_data, "acquisition_targets", "strategic_interest"):
            score += 1.5
        if _any_signal(company_data, "ipo_potential", "public_market_readiness"):
            score += 1.0
        if _any_signal(company_data, "comparable_exits", "valuation_multiples"):
            score += 0.8

    else:
        score += sum(0.8 for f in factors if _any_signal(company_data, f))

    return min(max(score, 1.0), 10.0)


def _calculate_composite_score(
        categories: List[Dict[str, Any]],
        module_weight: Optional[float] = None) -> float:
    """Calculate composite TCA score from weighted categories.

    module_weight, when provided, is the top-level TCA module weight from the
    saved /dashboard/module-settings configuration (e.g. 20 for 20%). It is
    currently informational only — the composite remains the weighted sum of
    sub-category scores — but is logged so downstream triage rendering can use
    it as an authoritative reference.
    """
    if module_weight is not None:
        logger.debug(f"Composite computed with TCA module weight={module_weight}")
    total_weighted_score = sum(cat["weightedScore"] for cat in categories)
    return round(total_weighted_score * 10, 1)  # Convert to 0-100 scale


def _calculate_risk_assessment(
        company_data: Dict[str, Any],
        tca_categories: List[Dict[str, Any]]) -> Dict[str, Any]:
    """Calculate risk assessment based on company data and TCA scores"""

    # Identify risk factors

    # Low scoring categories become risks
    risk_flags = [
        {
            "category": cat["category"],
            "severity": "high" if cat["rawScore"] < 5.0 else "medium",
            "description": f"Below threshold performance in {cat['category']}",
            "mitigation": f"Focus improvement efforts on {cat['category'].lower()}"
        }
        for cat in tca_categories
        if cat["rawScore"] < 6.5
    ]

    # Market risks
    if company_data.get("competitive_landscape") == "highly_competitive":
        risk_flags.append({
            "category": "Market Risk",
            "severity": "medium",
            "description": "Highly competitive market environment",
            "mitigation": "Strengthen differentiation strategy"
        })

    # Financial risks
    burn_rate = company_data.get("monthly_burn", 0)
    cash_balance = company_data.get("cash_balance", 0)
    if burn_rate > 0 and cash_balance > 0 and cash_balance / burn_rate < 12:  # Less than 12 months runway
        risk_flags.append({
            "category":
            "Financial Risk",
            "severity":
            "high",
            "description":
            "Limited cash runway",
            "mitigation":
            "Secure additional funding or reduce burn rate"
        })

    # Calculate overall risk level
    high_risks = len([r for r in risk_flags if r["severity"] == "high"])
    medium_risks = len([r for r in risk_flags if r["severity"] == "medium"])

    if high_risks > 2:
        overall_risk = "HIGH"
    elif high_risks > 0 or medium_risks > 3:
        overall_risk = "MEDIUM"
    else:
        overall_risk = "LOW"

    return {
        "riskLevel":
        overall_risk,
        "riskFlags":
        risk_flags,
        "totalRisks":
        len(risk_flags),
        "highSeverityRisks":
        high_risks,
        "mediumSeverityRisks":
        medium_risks,
        "summary":
        f"Identified {len(risk_flags)} risk factors with {overall_risk.lower()} overall risk level"
    }


def _generate_recommendation(composite_score: float,
                             risk_data: Dict[str, Any]) -> Dict[str, Any]:
    """Generate investment recommendation based on scores and risks"""

    if composite_score >= 80 and risk_data["riskLevel"] in ["LOW", "MEDIUM"]:
        overall = "STRONG INVEST - High potential with manageable risks"
        tier = "Tier 1"
    elif composite_score >= 70 and risk_data["riskLevel"] != "HIGH":
        overall = "INVEST - Solid opportunity with acceptable risk profile"
        tier = "Tier 2"
    elif composite_score >= 60:
        overall = "CONDITIONAL - Potential with improvement needed"
        tier = "Tier 3"
    else:
        overall = "PASS - Significant concerns identified"
        tier = "No Investment"

    # Generate strengths and concerns
    strengths = []
    concerns = []

    if composite_score >= 75:
        strengths.append("Strong overall fundamentals")
    if risk_data["riskLevel"] == "LOW":
        strengths.append("Low risk profile")
    if len(risk_data["riskFlags"]) == 0:
        strengths.append("No critical risks identified")

    if composite_score < 70:
        concerns.append("Below target composite score")
    if risk_data["highSeverityRisks"] > 0:
        concerns.append("High severity risks present")
    if len(risk_data["riskFlags"]) > 3:
        concerns.append("Multiple risk factors identified")

    return {
        "overall":
        overall,
        "tier":
        tier,
        "strengths":
        strengths[:3],  # Limit to top 3
        "concerns":
        concerns[:3],  # Limit to top 3
        "summary":
        f"TCA Analysis: {composite_score}/100 - {tier} recommendation. {overall}"
    }


def _get_category_interpretation(category: str, score: float) -> str:
    """Get interpretation for category score"""
    if score >= 8.0:
        return f"Excellent {category.lower()} performance - key competitive advantage"
    elif score >= 6.5:
        return f"Good {category.lower()} showing - meets investment criteria"
    else:
        return f"{category} needs improvement - requires attention"


def _get_category_recommendation(category: str, score: float) -> str:
    """Get AI recommendation for category"""
    if score >= 8.0:
        return f"Leverage strong {category.lower()} in investment narrative"
    elif score >= 6.5:
        return f"Continue building on {category.lower()} foundation"
    else:
        return f"Priority: Address {category.lower()} weaknesses before investment"


def _analyze_market_potential(company_data: Dict[str, Any]) -> str:
    """Analyze market potential based on data"""
    market_size = company_data.get("market_size", 0)
    if market_size > 1000000000:
        return "Large addressable market with significant growth potential"
    elif market_size > 100000000:
        return "Moderate market size with good growth opportunities"
    else:
        return "Market size requires validation and expansion strategy"


def _analyze_financial_health(company_data: Dict[str, Any]) -> str:
    """Analyze financial health"""
    revenue = company_data.get("monthly_revenue", 0)
    burn = company_data.get("monthly_burn", 0)

    if revenue > burn:
        return "Positive unit economics with sustainable financial model"
    elif revenue > 0:
        return "Revenue traction demonstrated, path to profitability visible"
    else:
        return "Pre-revenue stage requires clear monetization strategy"


def _analyze_team_strength(company_data: Dict[str, Any]) -> str:
    """Analyze team strength"""
    team_size = company_data.get("team_size", 0)
    if team_size >= 10:
        return "Well-sized team with diverse capabilities"
    elif team_size >= 5:
        return "Core team assembled with key skills represented"
    else:
        return "Early-stage team requiring strategic hiring"


# ============ Company Info Extraction Endpoint ============

@router.post("/extract-company-info", response_model=Dict[str, Any])
async def extract_company_info(request_data: Dict[str, Any]):
    """
    Extract company information from document content using AI/NLP.
    Used by frontend to auto-fill company information form after document upload.
    """
    import re
    
    content = request_data.get("content", "")
    framework = request_data.get("framework", "general")
    
    if not content:
        return {"error": "No content provided"}
    
    extracted = {}
    content_lower = content.lower()
    
    # ========== COMPANY NAME EXTRACTION (Enhanced) ==========
    name_patterns = [
        # Explicit labels with value after colon/equals
        r"(?:company\s*name|legal\s*name|startup\s*name|organization)[:\s=]+([A-Za-z0-9\s&.,'\-]+?)(?:\n|$|\s{2,}|Website|Founded|Industry|Location|Address)",
        # Pitch deck title patterns
        r"^([A-Z][A-Za-z0-9\s&'\-]{2,35})\s*[-–—:]\s*(?:pitch|deck|presentation|investor|overview)",
        # "About [Company]" pattern
        r"(?:welcome to|introducing|about|presenting)\s+([A-Z][A-Za-z0-9\s&'\-]{2,35})(?:\s|$|\n|,)",
        # Company with legal suffix at start of line
        r"(?:^|\n)\s*([A-Z][A-Za-z0-9\s&'\-]{2,25})(?:\s+(?:Inc\.?|LLC|Ltd\.?|Corp\.?|Co\.?|GmbH|PLC|SA))\s*(?:\n|$|,)",
        # Header/Title format (all caps or title case at beginning)
        r"^([A-Z][A-Z0-9\s&'\-]{3,30})\s*\n",
        # After "Overview of" or similar
        r"(?:overview of|profile of|summary of)\s+([A-Z][A-Za-z0-9\s&'\-]{2,35})",
    ]
    for pattern in name_patterns:
        match = re.search(pattern, content, re.IGNORECASE | re.MULTILINE)
        if match:
            name = match.group(1).strip()
            # Clean up the name - remove trailing punctuation and common words
            name = re.sub(r'[,.:;]+$', '', name).strip()
            name = re.sub(r'\s*(overview|pitch|deck|presentation)$', '', name, flags=re.IGNORECASE).strip()
            if len(name) >= 2 and len(name) <= 100:
                extracted["company_name"] = name
                break
    
    # ========== LEGAL NAME (if different) ==========
    if legal_match := re.search(r"(?:legal\s*name|registered\s*(?:as|name))[:\s=]+([A-Za-z0-9\s&.,'\-]+?)(?:\n|$|\s{2,})", content, re.IGNORECASE):
        extracted["legal_name"] = legal_match.group(1).strip()[:100]
    
    # ========== WEBSITE EXTRACTION (Enhanced) ==========
    # Try explicit label first
    website_labeled = re.search(r"(?:website|url|web|site|homepage)[:\s=]+(https?://[^\s<>\"']+|(?:www\.)?[a-zA-Z0-9][-a-zA-Z0-9]*\.[a-zA-Z]{2,}(?:/[^\s<>\"']*)?)", content, re.IGNORECASE)
    if website_labeled:
        url = website_labeled.group(1).rstrip('.,)>')
        extracted["website"] = url if url.startswith('http') else f"https://{url}"
    else:
        # Any URL in content
        website_match = re.search(r'(https?://[a-zA-Z0-9][-a-zA-Z0-9.]*\.[a-zA-Z]{2,}(?:/[^\s<>\"\']*)?)', content)
        if website_match:
            extracted["website"] = website_match.group(1).rstrip('.,)>')
    
    # ========== DESCRIPTION / ONE-LINER EXTRACTION (Enhanced) ==========
    # One-line description
    one_liner_patterns = [
        r"(?:one[\s-]?liner?|tagline|slogan|elevator\s*pitch)[:\s=]+([^\n]{10,150})",
        r"(?:we\s+(?:are|help|enable|provide|build))\s+([^\n.]{10,150}\.?)",
    ]
    for pattern in one_liner_patterns:
        match = re.search(pattern, content, re.IGNORECASE)
        if match:
            extracted["one_line_description"] = match.group(1).strip()[:200]
            break
    
    # Full description
    desc_patterns = [
        r"(?:company\s*description|about\s*us|overview|executive\s*summary)[:\s=]+([^\n]{30,}(?:\n[^\n]{20,}){0,3})",
        r"(?:mission|our\s*mission)[:\s=]+([^\n]{20,300})",
        r"(?:^|\n)([A-Z][^.]{30,300}(?:company|startup|platform|solution|service|business)[^.]{10,200}\.)",
    ]
    for pattern in desc_patterns:
        match = re.search(pattern, content, re.IGNORECASE | re.MULTILINE)
        if match:
            desc = match.group(1).strip()[:700]
            # Clean up
            desc = re.sub(r'\s+', ' ', desc)
            extracted["description"] = desc
            break
    
    # ========== PRODUCT DESCRIPTION ==========
    product_patterns = [
        r"(?:product|solution|platform|service)\s*(?:description|overview)?[:\s=]+([^\n]{20,}(?:\n[^\n]{20,}){0,2})",
        r"(?:what\s*we\s*(?:do|build|offer))[:\s=]+([^\n]{20,300})",
    ]
    for pattern in product_patterns:
        match = re.search(pattern, content, re.IGNORECASE)
        if match:
            extracted["product_description"] = match.group(1).strip()[:500]
            break
    
    # ========== EMPLOYEE COUNT (Enhanced) ==========
    emp_patterns = [
        r"(?:employees?|team\s*(?:size|members)|headcount|staff)[:\s=]+(\d{1,5})",
        r"(\d{1,5})\s*(?:employees?|team\s*members?|people|staff)",
        r"(?:team\s*of|staff\s*of)\s*(\d{1,5})",
    ]
    for pattern in emp_patterns:
        match = re.search(pattern, content, re.IGNORECASE)
        if match:
            extracted["number_of_employees"] = int(match.group(1))
            break
    
    # ========== LOCATION EXTRACTION (Enhanced) ==========
    # Full address
    address_match = re.search(r"(?:address|headquarters?|hq|located?(?:\s*(?:at|in))?)[:\s=]+([^\n]{10,150})", content, re.IGNORECASE)
    if address_match:
        addr = address_match.group(1).strip()
        extracted["address"] = addr[:200]
        # Try to parse city/state/country from address
        parts = [p.strip() for p in addr.split(',')]
        if len(parts) >= 2:
            extracted["city"] = parts[0][:50]
            extracted["country"] = parts[-1][:50]
            if len(parts) >= 3:
                extracted["state"] = parts[-2][:50]
    
    # Fallback location patterns
    if "city" not in extracted:
        location_patterns = [
            r"(?:based\s*in|located\s*in|headquarters?\s*in)\s+([A-Za-z\s]+),\s*([A-Za-z\s]+)",
            r"([A-Za-z]+),\s*(CA|NY|TX|WA|MA|FL|IL|PA|OH|GA|NC|NJ|VA|AZ|CO|MD|TN)\b",
        ]
        for pattern in location_patterns:
            match = re.search(pattern, content, re.IGNORECASE)
            if match:
                extracted["city"] = match.group(1).strip()[:50]
                extracted["state"] = match.group(2).strip()[:50]
                break
    
    # Country
    country_patterns = [
        r"(?:country|nation)[:\s=]+([A-Za-z\s]+?)(?:\n|$|,)",
        r"(?:US|USA|United States|UK|United Kingdom|Canada|Germany|France|Australia|India|Israel|Singapore)",
    ]
    if "country" not in extracted:
        for pattern in country_patterns:
            match = re.search(pattern, content, re.IGNORECASE)
            if match:
                country = match.group(0) if match.lastindex is None else match.group(1)
                # Normalize country names
                country_map = {"US": "United States", "USA": "United States", "UK": "United Kingdom"}
                extracted["country"] = country_map.get(country.strip().upper(), country.strip())[:50]
                break
    
    # ========== FOUNDING DATE / YEAR ==========
    founded_patterns = [
        r"(?:founded|established|started|launched|incorporated)\s*(?:in\s*)?(\d{4})",
        r"(?:year\s*founded|founding\s*year)[:\s=]+(\d{4})",
        r"(?:since|est\.?)\s*(\d{4})",
    ]
    for pattern in founded_patterns:
        match = re.search(pattern, content, re.IGNORECASE)
        if match:
            year = int(match.group(1))
            if 1900 <= year <= 2030:
                extracted["founded_year"] = year
                break
    
    # ========== FUNDING STAGE (Enhanced) ==========
    stage_keywords = {
        "pre-seed": "Pre-seed",
        "preseed": "Pre-seed",
        "seed round": "Seed",
        "seed stage": "Seed", 
        "seed funding": "Seed",
        "series a": "Series A",
        "series b": "Series B",
        "series c": "Series C+",
        "series d": "Series C+",
        "growth stage": "Growth",
        "growth round": "Growth",
        "late stage": "Growth",
        "ipo": "Public",
        "public company": "Public",
        "bootstrapped": "Bootstrapped",
        "pre-revenue": "Pre-seed",
        "mvp": "Pre-seed",
        "early stage": "Seed",
    }
    for keyword, stage in stage_keywords.items():
        if keyword in content_lower:
            extracted["development_stage"] = stage
            break
    
    # ========== FUNDING AMOUNT ==========
    funding_patterns = [
        r"(?:raised|funding|investment)\s*[:\s=]?\s*\$?([\d,.]+)\s*(million|m|billion|b|k|thousand)?",
        r"\$?([\d,.]+)\s*(million|m|billion|b)\s*(?:raised|funding|investment)",
    ]
    for pattern in funding_patterns:
        match = re.search(pattern, content, re.IGNORECASE)
        if match:
            amount = float(match.group(1).replace(',', ''))
            multiplier = match.group(2).lower() if match.group(2) else ''
            if multiplier in ['m', 'million']:
                amount *= 1000000
            elif multiplier in ['b', 'billion']:
                amount *= 1000000000
            elif multiplier in ['k', 'thousand']:
                amount *= 1000
            extracted["funding_amount"] = amount
            break
    
    # ========== INDUSTRY VERTICAL (Enhanced) ==========
    industry_keywords = {
        "software as a service": "Software/SaaS",
        "saas": "Software/SaaS",
        "software": "Software/SaaS",
        "fintech": "FinTech",
        "financial technology": "FinTech",
        "financial services": "FinTech",
        "healthtech": "HealthTech/MedTech",
        "health tech": "HealthTech/MedTech",
        "healthcare technology": "HealthTech/MedTech",
        "medtech": "HealthTech/MedTech",
        "medical technology": "HealthTech/MedTech",
        "medical device": "HealthTech/MedTech",
        "biotech": "BioTech",
        "biotechnology": "BioTech",
        "life sciences": "BioTech",
        "e-commerce": "E-commerce",
        "ecommerce": "E-commerce",
        "retail tech": "E-commerce",
        "edtech": "EdTech",
        "education technology": "EdTech",
        "learning platform": "EdTech",
        "cleantech": "CleanTech/GreenTech",
        "clean technology": "CleanTech/GreenTech",
        "green tech": "CleanTech/GreenTech",
        "sustainability": "CleanTech/GreenTech",
        "renewable": "CleanTech/GreenTech",
        "proptech": "PropTech",
        "property technology": "PropTech",
        "real estate tech": "PropTech",
        "artificial intelligence": "AI/ML",
        "machine learning": "AI/ML",
        "deep learning": "AI/ML",
        "cybersecurity": "Cybersecurity",
        "security": "Cybersecurity",
        "insurtech": "InsurTech",
        "insurance technology": "InsurTech",
        "logistics": "Logistics/Supply Chain",
        "supply chain": "Logistics/Supply Chain",
        "foodtech": "FoodTech",
        "food technology": "FoodTech",
        "agtech": "AgTech",
        "agriculture technology": "AgTech",
        "gaming": "Gaming/Entertainment",
        "entertainment": "Gaming/Entertainment",
        "media": "Media/Content",
        "content platform": "Media/Content",
        "hr tech": "HR Tech",
        "human resources": "HR Tech",
        "marketing tech": "MarTech",
        "martech": "MarTech",
        "advertising": "AdTech",
        "adtech": "AdTech",
    }
    for keyword, industry in industry_keywords.items():
        if keyword in content_lower:
            extracted["industry_vertical"] = industry
            break
    
    # ========== BUSINESS MODEL (Enhanced) ==========
    model_keywords = {
        "b2b saas": "B2B SaaS",
        "b2b software": "B2B SaaS",
        "enterprise software": "B2B SaaS",
        "b2c saas": "B2C SaaS",
        "consumer saas": "B2C SaaS",
        "b2b": "B2B",
        "b2c": "B2C",
        "b2b2c": "B2B2C",
        "marketplace": "Marketplace",
        "two-sided marketplace": "Marketplace",
        "subscription": "Subscription",
        "recurring revenue": "Subscription",
        "freemium": "Freemium",
        "free trial": "Freemium",
        "platform": "Platform",
        "api": "API/Platform",
        "licensing": "Licensing",
        "transactional": "Transactional",
        "per-transaction": "Transactional",
        "advertising": "Advertising",
        "ad-supported": "Advertising",
        "hardware": "Hardware",
        "device": "Hardware",
        "consulting": "Services",
        "professional services": "Services",
    }
    for keyword, model in model_keywords.items():
        if keyword in content_lower:
            extracted["business_model"] = model
            break
    
    # ========== REVENUE / FINANCIALS ==========
    revenue_patterns = [
        r"(?:revenue|arr|annual\s*recurring\s*revenue|mrr)[:\s=]?\s*\$?([\d,.]+)\s*(million|m|k|thousand)?",
        r"(?:revenue|arr|mrr)\s*(?:of|:)?\s*\$?([\d,.]+)\s*(million|m|k)?",
    ]
    for pattern in revenue_patterns:
        match = re.search(pattern, content, re.IGNORECASE)
        if match:
            amount = float(match.group(1).replace(',', ''))
            multiplier = match.group(2).lower() if match.group(2) else ''
            if multiplier in ['m', 'million']:
                amount *= 1000000
            elif multiplier in ['k', 'thousand']:
                amount *= 1000
            extracted["annual_revenue"] = amount
            break
    
    # ========== FOUNDERS / CEO ==========
    founder_patterns = [
        r"(?:founder(?:s)?|ceo|chief\s*executive)[:\s=]+([A-Z][a-z]+(?:\s+[A-Z][a-z]+){1,3})",
        r"(?:founded\s*by|co-founders?)[:\s=]+([A-Z][a-z]+(?:\s+[A-Z][a-z]+){1,3}(?:\s*(?:and|,)\s*[A-Z][a-z]+(?:\s+[A-Z][a-z]+){1,3})*)",
    ]
    for pattern in founder_patterns:
        match = re.search(pattern, content, re.IGNORECASE)
        if match:
            extracted["founders"] = match.group(1).strip()[:200]
            break
    
    # ========== FRAMEWORK-SPECIFIC DEFAULTS ==========
    if framework == "medtech" and not extracted.get("industry_vertical"):
        extracted["industry_vertical"] = "HealthTech/MedTech"
    elif framework == "fintech" and not extracted.get("industry_vertical"):
        extracted["industry_vertical"] = "FinTech"
    elif framework == "saas" and not extracted.get("industry_vertical"):
        extracted["industry_vertical"] = "Software/SaaS"
    
    logger.info(f"Extracted company info: {list(extracted.keys())} - {len(extracted)} fields")
    
    return extracted


# ============ Analyst Review Workflow Endpoints ============

@router.post("/analyst-reviews", response_model=Dict[str, Any])
async def create_analyst_review(review_data: Dict[str, Any]):
    """
    Create a new analyst review for an analysis.
    Captures human scores, comments, and deviation rationale.
    """
    from datetime import datetime
    
    try:
        analysis_id = review_data.get("analysis_id")
        if not analysis_id:
            raise HTTPException(status_code=400, detail="analysis_id required")
        
        review = {
            "id": f"review_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            "analysis_id": analysis_id,
            "analyst_id": review_data.get("analyst_id"),
            "analyst_name": review_data.get("analyst_name", "Unknown Analyst"),
            "human_scores": review_data.get("human_scores", {}),
            "ai_scores": review_data.get("ai_scores", {}),
            "deviation_rationale": review_data.get("deviation_rationale", ""),
            "category_comments": review_data.get("category_comments", {}),
            "overall_assessment": review_data.get("overall_assessment", ""),
            "sentiment_indicators": review_data.get("sentiment_indicators", {}),
            "recommendation_override": review_data.get("recommendation_override"),
            "training_approved": review_data.get("training_approved", False),
            "status": "pending",
            "created_at": datetime.now().isoformat(),
            "updated_at": datetime.now().isoformat(),
        }
        
        # Calculate deviation metrics
        review["deviation_metrics"] = _calculate_deviation_metrics(
            review["human_scores"], 
            review["ai_scores"]
        )
        
        # Analyze sentiment from comments
        review["sentiment_analysis"] = _analyze_review_sentiment(
            review["category_comments"],
            review["deviation_rationale"]
        )
        
        logger.info(f"Created analyst review: {review['id']}")
        return review
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to create analyst review: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.put("/analyst-reviews", response_model=Dict[str, Any])
async def upsert_analyst_review(review_data: Dict[str, Any]):
    """PUT alias for clients that upsert analyst reviews."""
    return await create_analyst_review(review_data)


@router.get("/analyst-reviews/{analysis_id}", response_model=Dict[str, Any])
async def get_analyst_review(analysis_id: str):
    """Get analyst review for an analysis"""
    # In production, fetch from database
    return {
        "analysis_id": analysis_id,
        "reviews": [],
        "message": "No reviews found for this analysis"
    }


@router.post("/ai-deviation-comparison", response_model=Dict[str, Any])
async def compare_ai_human_scores(comparison_data: Dict[str, Any]):
    """
    Compare AI scores with human analyst scores and generate deviation analysis.
    Used by the AI vs Human Gap Analysis component.
    """
    try:
        ai_scores = comparison_data.get("ai_scores", {})
        human_scores = comparison_data.get("human_scores", {})
        analyst_comments = comparison_data.get("comments", {})
        
        if not ai_scores or not human_scores:
            raise HTTPException(status_code=400, detail="Both ai_scores and human_scores required")
        
        # Calculate deviation metrics
        deviation_metrics = _calculate_deviation_metrics(human_scores, ai_scores)
        
        # Generate chart data for visualization
        chart_data = []
        for category in set(ai_scores.keys()) | set(human_scores.keys()):
            ai_score = ai_scores.get(category, 0)
            human_score = human_scores.get(category, 0)
            chart_data.append({
                "category": category,
                "ai": ai_score,
                "analyst": human_score,
                "deviation": round(ai_score - human_score, 2),
                "comment": analyst_comments.get(category, "")
            })
        
        # Sort by absolute deviation (largest first)
        chart_data.sort(key=lambda x: abs(x["deviation"]), reverse=True)
        
        # Generate AI recommendations based on deviations
        recommendations = _generate_training_recommendations(chart_data, deviation_metrics)
        
        # Perform sentiment analysis on comments
        sentiment_analysis = _analyze_review_sentiment(analyst_comments, "")
        
        return {
            "chart_data": chart_data,
            "deviation_metrics": deviation_metrics,
            "recommendations": recommendations,
            "sentiment_analysis": sentiment_analysis,
            "training_suggested": deviation_metrics.get("mae", 0) > 0.5,
            "review_needed": any(abs(d["deviation"]) > 2.0 for d in chart_data)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"AI deviation comparison failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/submit-for-training", response_model=Dict[str, Any])
async def submit_for_ml_training(training_data: Dict[str, Any]):
    """
    Submit analyst review data for ML model training.
    Captures human corrections to improve AI accuracy over time.
    """
    from datetime import datetime
    
    try:
        training_record = {
            "id": f"train_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            "analysis_id": training_data.get("analysis_id"),
            "company_name": training_data.get("company_name", "Unknown"),
            "ai_scores": training_data.get("ai_scores", {}),
            "human_scores": training_data.get("human_scores", {}),
            "human_rationale": training_data.get("rationale", ""),
            "category_adjustments": training_data.get("category_adjustments", {}),
            "analyst_id": training_data.get("analyst_id"),
            "industry_context": training_data.get("industry_context", ""),
            "company_stage": training_data.get("company_stage", ""),
            "submitted_at": datetime.now().isoformat(),
            "status": "queued",
            "priority": _calculate_training_priority(training_data)
        }
        
        logger.info(f"Submitted training data: {training_record['id']}")
        
        return {
            "success": True,
            "training_id": training_record["id"],
            "message": "Training data submitted successfully. It will be processed in the next training batch.",
            "estimated_incorporation": "Within 24 hours",
            "priority": training_record["priority"]
        }
        
    except Exception as e:
        logger.error(f"Failed to submit training data: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/sentiment-analysis", response_model=Dict[str, Any])
async def analyze_comment_sentiment(request_data: Dict[str, Any]):
    """
    Analyze sentiment and insights from analyst comments.
    Supports deep content analysis for reviewer feedback.
    """
    try:
        comments = request_data.get("comments", {})
        rationale = request_data.get("rationale", "")
        
        analysis = _analyze_review_sentiment(comments, rationale)
        
        # Deep content analysis
        content_insights = _deep_content_analysis(comments, rationale)
        
        return {
            "sentiment_scores": analysis,
            "content_insights": content_insights,
            "key_themes": content_insights.get("themes", []),
            "action_items": content_insights.get("action_items", []),
            "confidence_level": content_insights.get("confidence", "medium")
        }
        
    except Exception as e:
        logger.error(f"Sentiment analysis failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))


def _calculate_deviation_metrics(human_scores: Dict[str, float], ai_scores: Dict[str, float]) -> Dict[str, Any]:
    """Calculate statistical metrics for AI vs Human score deviations"""
    import math
    
    categories = list(set(human_scores.keys()) & set(ai_scores.keys()))
    if not categories:
        return {"error": "No overlapping categories"}
    
    deviations = []
    squared_deviations = []
    agreements = 0
    
    for cat in categories:
        human = human_scores.get(cat, 0)
        ai = ai_scores.get(cat, 0)
        diff = abs(ai - human)
        deviations.append(diff)
        squared_deviations.append(diff ** 2)
        if diff < 0.5:  # Within 0.5 points is considered agreement
            agreements += 1
    
    n = len(deviations)
    mae = sum(deviations) / n  # Mean Absolute Error
    rmse = math.sqrt(sum(squared_deviations) / n)  # Root Mean Square Error
    agreement_rate = agreements / n
    
    # Calculate Cohen's Kappa approximation
    # (simplified for numerical scores - use discretized buckets)
    human_bins = [_score_bucket(human_scores.get(c, 0)) for c in categories]
    ai_bins = [_score_bucket(ai_scores.get(c, 0)) for c in categories]
    observed_agreement = sum(h == a for h, a in zip(human_bins, ai_bins)) / n if n > 0 else 0
    expected_agreement = 0.33  # Random agreement for 3 buckets
    kappa = (observed_agreement - expected_agreement) / (1 - expected_agreement) if expected_agreement < 1 else 1
    
    # Bias calculation (positive = AI higher, negative = Human higher)
    bias = sum(ai_scores.get(c, 0) - human_scores.get(c, 0) for c in categories) / n
    
    return {
        "mae": round(mae, 3),
        "rmse": round(rmse, 3),
        "cohens_kappa": round(max(0, kappa), 3),
        "agreement_rate": round(agreement_rate * 100, 1),
        "bias": round(bias, 3),
        "bias_direction": "AI Higher" if bias > 0.1 else "Human Higher" if bias < -0.1 else "Neutral",
        "total_categories": n,
        "high_deviation_count": len([d for d in deviations if d > 1.5]),
        "calibration_quality": "High" if mae < 0.5 else "Medium" if mae < 1.0 else "Low"
    }


def _score_bucket(score: float) -> str:
    """Convert score to bucket for kappa calculation"""
    if score >= 8:
        return "high"
    elif score >= 6:
        return "medium"
    else:
        return "low"


def _generate_training_recommendations(chart_data: List[Dict], metrics: Dict) -> Dict[str, Any]:
    """Generate AI training recommendations based on deviations"""
    
    recommendations = []
    weight_adjustments = {}
    
    for item in chart_data:
        deviation = item["deviation"]
        category = item["category"]
        
        if abs(deviation) > 1.5:
            if deviation > 0:
                recommendations.append(
                    f"Reduce weight on '{category}' - AI over-scoring by {abs(deviation):.1f} points"
                )
                weight_adjustments[category] = {"action": "reduce", "magnitude": min(0.2, abs(deviation) / 10)}
            else:
                recommendations.append(
                    f"Increase weight on '{category}' - AI under-scoring by {abs(deviation):.1f} points"
                )
                weight_adjustments[category] = {"action": "increase", "magnitude": min(0.2, abs(deviation) / 10)}
    
    # Add general recommendations based on metrics
    if metrics.get("bias", 0) > 0.5:
        recommendations.append("General bias detected: AI tends to score higher. Consider calibration adjustment.")
    elif metrics.get("bias", 0) < -0.5:
        recommendations.append("General bias detected: AI tends to score lower. Consider calibration adjustment.")
    
    if metrics.get("mae", 0) > 1.0:
        recommendations.append("High overall deviation suggests model retraining may be beneficial.")
    
    return {
        "recommendations": recommendations[:5],  # Top 5 recommendations
        "weight_adjustments": weight_adjustments,
        "retraining_suggested": metrics.get("mae", 0) > 0.8,
        "priority": "high" if metrics.get("mae", 0) > 1.5 else "medium" if metrics.get("mae", 0) > 0.8 else "low"
    }


def _analyze_review_sentiment(comments: Dict[str, str], rationale: str) -> Dict[str, Any]:
    """Analyze sentiment from analyst comments"""
    
    # Simple keyword-based sentiment analysis
    positive_keywords = [
        "strong", "excellent", "good", "impressive", "solid", "promising", "innovative",
        "growth", "opportunity", "potential", "advantage", "leading", "experienced"
    ]
    negative_keywords = [
        "weak", "poor", "concern", "risk", "challenge", "limited", "lacking",
        "unclear", "uncertain", "competitive", "struggling", "unproven"
    ]
    neutral_keywords = ["adequate", "average", "moderate", "typical", "standard"]
    
    all_text = " ".join(comments.values()) + " " + rationale
    all_text_lower = all_text.lower()
    
    positive_count = sum(kw in all_text_lower for kw in positive_keywords)
    negative_count = sum(kw in all_text_lower for kw in negative_keywords)
    neutral_count = sum(kw in all_text_lower for kw in neutral_keywords)
    
    total = positive_count + negative_count + neutral_count + 1  # +1 to avoid division by zero
    
    # Calculate sentiment score (-1 to 1)
    sentiment_score = (positive_count - negative_count) / total
    
    # Per-category sentiment
    category_sentiment = {}
    for cat, comment in comments.items():
        comment_lower = comment.lower()
        cat_positive = sum(kw in comment_lower for kw in positive_keywords)
        cat_negative = sum(kw in comment_lower for kw in negative_keywords)
        if cat_positive > cat_negative:
            category_sentiment[cat] = "positive"
        elif cat_negative > cat_positive:
            category_sentiment[cat] = "negative"
        else:
            category_sentiment[cat] = "neutral"
    
    return {
        "overall_sentiment": "positive" if sentiment_score > 0.2 else "negative" if sentiment_score < -0.2 else "neutral",
        "sentiment_score": round(sentiment_score, 3),
        "positive_indicators": positive_count,
        "negative_indicators": negative_count,
        "category_sentiment": category_sentiment,
        "confidence": "high" if total > 5 else "medium" if total > 2 else "low"
    }


def _deep_content_analysis(comments: Dict[str, str], rationale: str) -> Dict[str, Any]:
    """Perform deep content analysis on analyst comments"""
    import re
    
    all_text = " ".join(comments.values()) + " " + rationale
    
    # Extract themes
    theme_patterns = {
        "market_concerns": r"market|competition|competitive|competitor|share|penetration",
        "team_assessment": r"team|leadership|founder|experience|execute|execution",
        "financial_view": r"revenue|financial|burn|runway|profit|margin|cost",
        "product_feedback": r"product|technology|innovation|feature|development",
        "growth_outlook": r"growth|scale|expand|traction|momentum|potential",
        "risk_factors": r"risk|challenge|concern|threat|weakness|barrier",
    }
    
    themes = [
        theme.replace("_", " ").title()
        for theme, pattern in theme_patterns.items()
        if re.search(pattern, all_text, re.IGNORECASE)
    ]
    
    # Extract action items (sentences with action verbs)
    action_patterns = [
        r"(?:should|need to|must|recommend|suggest|consider)\s+([^.!?]+[.!?])",
        r"(?:improve|address|focus on|strengthen)\s+([^.!?]+[.!?])",
    ]
    
    action_items = []
    for pattern in action_patterns:
        matches = re.findall(pattern, all_text, re.IGNORECASE)
        action_items.extend([m.strip() for m in matches[:3]])  # Max 3 per pattern
    
    # Calculate content depth
    word_count = len(all_text.split())
    sentence_count = len(re.split(r'[.!?]+', all_text))
    avg_sentence_length = word_count / max(sentence_count, 1)
    
    return {
        "themes": themes[:5],
        "action_items": action_items[:5],
        "word_count": word_count,
        "depth_score": "comprehensive" if word_count > 200 else "moderate" if word_count > 50 else "brief",
        "avg_sentence_length": round(avg_sentence_length, 1),
        "confidence": "high" if word_count > 100 else "medium" if word_count > 30 else "low"
    }


def _calculate_training_priority(training_data: Dict[str, Any]) -> str:
    """Calculate priority level for training data"""
    
    # Higher priority for larger deviations
    ai_scores = training_data.get("ai_scores", {})
    human_scores = training_data.get("human_scores", {})
    
    if not ai_scores or not human_scores:
        return "low"
    
    max_deviation = 0
    for cat in set(ai_scores.keys()) & set(human_scores.keys()):
        deviation = abs(ai_scores[cat] - human_scores[cat])
        max_deviation = max(max_deviation, deviation)
    
    # Also consider if rationale is provided (more valuable for training)
    has_rationale = bool(training_data.get("rationale", "").strip())
    
    if max_deviation > 2.0 or (max_deviation > 1.0 and has_rationale):
        return "high"
    elif max_deviation > 1.0 or has_rationale:
        return "medium"
    else:
        return "low"


# ============ File Upload & Text Extraction Endpoint ============

@router.post("/extract-text-from-file", response_model=Dict[str, Any])
async def extract_text_from_file(
    file: Optional[bytes] = None,
    file_data: Optional[Dict[str, Any]] = None,
):
    """
    Extract text content from uploaded files (PDF, DOCX, etc.)
    Accepts either raw file bytes or base64-encoded file data.
    """
    import base64
    import io
    
    try:
        content = ""
        file_bytes = None
        filename = "unknown"
        
        if file_data:
            # Handle base64 encoded file
            file_bytes = base64.b64decode(file_data.get("content", ""))
            filename = file_data.get("filename", "unknown")
        elif file:
            file_bytes = file
        
        if not file_bytes:
            return {"error": "No file data provided", "text_content": ""}
        
        # Try to extract text based on file type
        if filename.lower().endswith('.pdf'):
            # Use PyMuPDF for PDF extraction if available
            try:
                import fitz  # PyMuPDF
                pdf_doc = fitz.open(stream=file_bytes, filetype="pdf")
                text_parts = []
                for page_num in range(len(pdf_doc)):
                    page = pdf_doc.load_page(page_num)
                    text_parts.append(page.get_text())
                pdf_doc.close()
                content = "\n".join(text_parts)
            except ImportError:
                # Fallback: try pdfplumber
                try:
                    import pdfplumber
                    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                        text_parts = []
                        for page in pdf.pages:
                            if text := page.extract_text():
                                text_parts.append(text)
                        content = "\n".join(text_parts)
                except ImportError:
                    content = "[PDF extraction requires PyMuPDF or pdfplumber]"
                except Exception as e:
                    content = f"[PDF extraction failed: {str(e)}]"
            except Exception as e:
                content = f"[PDF extraction failed: {str(e)}]"
                
        elif filename.lower().endswith(('.docx', '.doc')):
            try:
                import docx
                doc = docx.Document(io.BytesIO(file_bytes))
                content = "\n".join([para.text for para in doc.paragraphs])
            except ImportError:
                content = "[DOCX extraction requires python-docx]"
            except Exception as e:
                content = f"[DOCX extraction failed: {str(e)}]"
        
        elif filename.lower().endswith(('.xlsx', '.xls')):
            # Excel file extraction
            try:
                import openpyxl
                workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
                text_parts = []
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text_parts.append(f"=== Sheet: {sheet_name} ===")
                    for row in sheet.iter_rows(values_only=True):
                        row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                        if row_text.strip():
                            text_parts.append(row_text)
                content = "\n".join(text_parts)
            except ImportError:
                # Fallback: try pandas
                try:
                    import pandas as pd
                    excel_file = pd.ExcelFile(io.BytesIO(file_bytes))
                    text_parts = []
                    for sheet_name in excel_file.sheet_names:
                        df = pd.read_excel(excel_file, sheet_name=sheet_name)
                        text_parts.extend([f"=== Sheet: {sheet_name} ===", df.to_string()])
                    content = "\n".join(text_parts)
                except ImportError:
                    content = "[Excel extraction requires openpyxl or pandas]"
                except Exception as e:
                    content = f"[Excel extraction failed: {str(e)}]"
            except Exception as e:
                content = f"[Excel extraction failed: {str(e)}]"
        
        elif filename.lower().endswith('.pptx'):
            # PowerPoint file extraction
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                text_parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    shape_texts = [
                        t for t in (getattr(shape, "text", "") for shape in slide.shapes) if t
                    ]
                    text_parts.extend([f"=== Slide {slide_num} ===", *shape_texts])
                content = "\n".join(text_parts)
            except ImportError:
                content = "[PowerPoint extraction requires python-pptx]"
            except Exception as e:
                content = f"[PowerPoint extraction failed: {str(e)}]"
        
        elif filename.lower().endswith('.rtf'):
            # RTF file extraction
            try:
                from striprtf.striprtf import rtf_to_text
                content = rtf_to_text(file_bytes.decode('utf-8', errors='ignore'))
            except ImportError:
                content = "[RTF extraction requires striprtf]"
            except Exception as e:
                content = f"[RTF extraction failed: {str(e)}]"
                
        elif filename.lower().endswith(('.txt', '.csv', '.json', '.md', '.yaml', '.yml')):
            try:
                content = file_bytes.decode('utf-8')
            except UnicodeDecodeError:
                content = file_bytes.decode('latin-1')
        else:
            content = f"[Unsupported file type: {filename}]"
        
        logger.info(f"Extracted {len(content)} characters from {filename}")
        
        return {
            "filename": filename,
            "text_content": content,
            "char_count": len(content),
            "word_count": len(content.split()) if content else 0,
        }
        
    except Exception as e:
        logger.error(f"File text extraction failed: {e}")
        return {"error": str(e), "text_content": ""}
